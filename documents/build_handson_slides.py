"""Build the CSHA modularity hands-on slide deck from generated figures.

Run the numbered analyses first, then build the deck with:

    poetry run python documents/build_handson_slides.py

The source talk ``CSHA_082426.pptx`` is deliberately not opened or modified by
this script.  It remains an ignored local reference; this builder creates a
small, course-specific deck that can be regenerated from repository outputs.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "results" / "figures"
OUTPUT = ROOT / "documents" / "CSHA_handson_tutorial.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 34, 54)
INK = RGBColor(30, 43, 55)
SLATE = RGBColor(88, 103, 116)
MUTED = RGBColor(118, 132, 143)
TEAL = RGBColor(0, 151, 157)
BLUE = RGBColor(62, 106, 181)
CORAL = RGBColor(224, 82, 77)
GOLD = RGBColor(224, 164, 54)
PALE = RGBColor(244, 247, 249)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(213, 222, 228)

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"
FONT_MONO = "Aptos Mono"


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, size=18, color=INK, gap=7):
    """Add separate paragraphs specified as (text, bold, color) tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.word_wrap = True
    for index, (text, is_bold, line_color) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = line_color or color
    return box


def add_footer(slide, step=None):
    add_rect(slide, Inches(0.55), Inches(7.12), Inches(12.23), Inches(0.012), LINE)
    add_text(
        slide,
        "CSHA Neural Data Science · Functional-network modularity",
        Inches(0.58),
        Inches(7.19),
        Inches(6.8),
        Inches(0.18),
        size=8,
        color=MUTED,
    )
    if step:
        add_text(
            slide,
            step,
            Inches(10.1),
            Inches(7.19),
            Inches(2.65),
            Inches(0.18),
            size=8,
            color=MUTED,
            align=PP_ALIGN.RIGHT,
        )


def add_header(slide, title, *, eyebrow=None, step=None):
    if eyebrow:
        add_text(
            slide,
            eyebrow.upper(),
            Inches(0.62),
            Inches(0.30),
            Inches(7.0),
            Inches(0.24),
            size=9,
            color=TEAL,
            bold=True,
        )
    add_text(
        slide,
        title,
        Inches(0.60),
        Inches(0.60 if eyebrow else 0.38),
        Inches(12.1),
        Inches(0.82),
        size=28,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_footer(slide, step)


def add_slide(prs, title, *, eyebrow=None, step=None, background=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, background)
    add_header(slide, title, eyebrow=eyebrow, step=step)
    return slide


def add_picture_contain(slide, path, x, y, w, h, *, border=True):
    path = Path(path)
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio >= box_ratio:
        pic_w = w
        pic_h = int(w / image_ratio)
    else:
        pic_h = h
        pic_w = int(h * image_ratio)
    pic_x = int(x + (w - pic_w) / 2)
    pic_y = int(y + (h - pic_h) / 2)
    if border:
        add_rect(slide, x, y, w, h, WHITE, LINE, radius=True)
    return slide.shapes.add_picture(str(path), pic_x, pic_y, pic_w, pic_h)


def add_caption(slide, text, x, y, w, *, align=PP_ALIGN.LEFT):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        Inches(0.32),
        size=9,
        color=SLATE,
        align=align,
    )


def add_pill(slide, text, x, y, w, *, color=TEAL):
    add_rect(slide, x, y, w, Inches(0.36), color, color, radius=True)
    add_text(
        slide,
        text,
        x,
        y + Inches(0.02),
        w,
        Inches(0.25),
        size=10,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_takeaway(slide, text, x, y, w, *, color=TEAL, height=0.72, size=13):
    h = Inches(height)
    add_rect(slide, x, y, w, h, PALE, PALE, radius=True)
    add_rect(slide, x, y, Inches(0.07), h, color, color, radius=True)
    add_text(
        slide,
        text,
        x + Inches(0.18),
        y + Inches(0.09),
        w - Inches(0.30),
        h - Inches(0.18),
        size=size,
        color=INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_script_badge(slide, number, label, x, y, *, color=TEAL):
    add_rect(slide, x, y, Inches(0.52), Inches(0.52), color, color, radius=True)
    add_text(
        slide,
        number,
        x,
        y + Inches(0.01),
        Inches(0.52),
        Inches(0.38),
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        label,
        x + Inches(0.67),
        y + Inches(0.05),
        Inches(3.15),
        Inches(0.40),
        size=13,
        color=NAVY,
        bold=True,
    )


def figure_file(name, *fallback_names):
    """Return a generated figure, allowing a legacy filename during migration."""
    for candidate in (name, *fallback_names):
        path = FIGURES / candidate
        if path.exists():
            return path
    return FIGURES / name


def require_figures():
    names = [
        "01_raw_dff_traces.png",
        "01_spatial_distribution_by_area.png",
        "01_eeg_emg_state_classification.png",
        "02_connectivity_matrices.png",
        "02_fixed_density_graphs.png",
        "03_modularity_pipeline_sleep.png",
        "03_spatial_modules_sleep.png",
        "04_sample_state_comparison.png",
        "05_sample_coarse_grain_modularity.png",
        "06_sample_module_spatial_distribution.png",
        "07_all_mice_state_comparison.png",
        "07_modularity_per_mouse_sleep.png",
        "07_modularity_per_mouse_ane.png",
        "08_all_mice_coarse_grain_modularity.png",
        "08_all_mice_coarse_grain_modules.png",
        "09_all_mice_module_maps.png",
        "09_all_mice_same_module_vs_distance.png",
    ]
    missing = [name for name in names if not (FIGURES / name).exists()]
    movie_overview = figure_file(
        "multiscale_overview_mouse02_sleep_nrem.png",
        "07_multiscale_overview_mouse02_sleep_nrem.png",
    )
    if not movie_overview.exists():
        missing.append("multiscale_overview_mouse02_sleep_nrem.png")
    if missing:
        formatted = "\n  - ".join(missing)
        raise FileNotFoundError(
            "Run the numbered scripts before building the slides. Missing:\n  - "
            f"{formatted}"
        )


def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(0.16), SLIDE_H, TEAL)
    add_text(
        slide,
        "CSHA · NEURAL DATA SCIENCE",
        Inches(0.82),
        Inches(0.72),
        Inches(5.5),
        Inches(0.30),
        size=11,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "Functional-network\nmodularity across brain states",
        Inches(0.78),
        Inches(1.42),
        Inches(7.4),
        Inches(1.80),
        size=34,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
        line_spacing=0.9,
    )
    add_text(
        slide,
        "A hands-on path from one complete calcium-imaging recording\n"
        "to a reproducible all-mice analysis",
        Inches(0.82),
        Inches(3.58),
        Inches(6.5),
        Inches(0.9),
        size=18,
        color=RGBColor(207, 220, 229),
        line_spacing=1.1,
    )
    for index, (number, label, color) in enumerate(
        [
            ("A", "mouse02_sleep · scripts 00–06", TEAL),
            ("B", "all mice · scripts 07–09", GOLD),
        ]
    ):
        y = Inches(5.10 + 0.72 * index)
        add_rect(slide, Inches(0.82), y, Inches(0.46), Inches(0.46), color, color, radius=True)
        add_text(
            slide,
            number,
            Inches(0.82),
            y + Inches(0.015),
            Inches(0.46),
            Inches(0.32),
            size=13,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            Inches(1.48),
            y + Inches(0.05),
            Inches(4.8),
            Inches(0.32),
            size=15,
            color=WHITE,
            bold=True,
        )
    add_picture_contain(
        slide,
        FIGURES / "03_spatial_modules_sleep.png",
        Inches(8.2),
        Inches(0.7),
        Inches(4.45),
        Inches(5.95),
        border=False,
    )
    add_text(
        slide,
        "Kiyooka & Oomoto et al. · Cell Reports (2026)",
        Inches(8.0),
        Inches(6.75),
        Inches(4.7),
        Inches(0.25),
        size=9,
        color=RGBColor(166, 186, 199),
        align=PP_ALIGN.RIGHT,
    )


def build_roadmap(prs):
    slide = add_slide(prs, "Start small, then test robustness", eyebrow="Tutorial roadmap")
    add_text(
        slide,
        "Everyone completes Track A. Track B is the launch point for independent projects.",
        Inches(0.63),
        Inches(1.19),
        Inches(11.8),
        Inches(0.40),
        size=16,
        color=SLATE,
    )

    tracks = [
        (
            "TRACK A · RECOMMENDED",
            "One complete recording",
            "mouse02_sleep + synchronized EEG/EMG",
            "00  download\n01  inspect signals + physiology\n02  build equal-density networks\n03  identify modules\n04  compare states in one mouse\n05–06  explore spatial scale\noptional  multiscale movie",
            TEAL,
        ),
        (
            "TRACK B · RESEARCH EXTENSION",
            "All biological mice",
            "5 sleep mice + 4 anesthesia mice",
            "00 --all  download the cohort\n07  compare modularity across mice\n08  test the scale dependence\n09  quantify module geography\n\nSet PAPER_MODE = True for the full run.",
            GOLD,
        ),
    ]
    for index, (tag, title, subtitle, body, color) in enumerate(tracks):
        x = Inches(0.63 + 6.18 * index)
        add_rect(slide, x, Inches(1.77), Inches(5.86), Inches(4.77), PALE, LINE, radius=True)
        add_pill(slide, tag, x + Inches(0.28), Inches(2.04), Inches(2.55), color=color)
        add_text(
            slide,
            title,
            x + Inches(0.28),
            Inches(2.63),
            Inches(5.20),
            Inches(0.46),
            size=24,
            color=NAVY,
            bold=True,
            font=FONT_DISPLAY,
        )
        add_text(
            slide,
            subtitle,
            x + Inches(0.28),
            Inches(3.17),
            Inches(5.20),
            Inches(0.36),
            size=13,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + Inches(0.34),
            Inches(3.72),
            Inches(4.95),
            Inches(2.30),
            size=14,
            color=INK,
            font=FONT_MONO,
            line_spacing=1.16,
        )


def build_download(prs):
    slide = add_slide(
        prs,
        "00 · Download one real experiment",
        eyebrow="Track A · complete example",
        step="script 00",
    )
    add_text(
        slide,
        "The default download provides one complete calcium + physiology session.",
        Inches(0.63),
        Inches(1.19),
        Inches(8.0),
        Inches(0.40),
        size=16,
        color=SLATE,
    )
    files = [
        ("CALCIUM", "mouse02_sleep.mat", "~1.09 GB", "6,574 neurons · 18,700 frames", TEAL),
        (
            "PHYSIOLOGY",
            "mouse02_sleep_physiological_data.mat",
            "~0.36 GB",
            "5-kHz EEG + EMG · frame triggers",
            BLUE,
        ),
    ]
    for index, (tag, name, size, details, color) in enumerate(files):
        x = Inches(0.64 + 6.18 * index)
        add_rect(slide, x, Inches(1.86), Inches(5.84), Inches(2.28), WHITE, LINE, radius=True)
        add_pill(slide, tag, x + Inches(0.27), Inches(2.10), Inches(1.24), color=color)
        add_text(
            slide,
            name,
            x + Inches(0.27),
            Inches(2.70),
            Inches(5.24),
            Inches(0.43),
            size=17,
            color=NAVY,
            bold=True,
            font=FONT_MONO,
        )
        add_text(
            slide,
            size,
            x + Inches(0.27),
            Inches(3.25),
            Inches(1.25),
            Inches(0.32),
            size=16,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            details,
            x + Inches(1.55),
            Inches(3.27),
            Inches(3.9),
            Inches(0.32),
            size=13,
            color=SLATE,
        )
    add_rect(slide, Inches(0.64), Inches(4.58), Inches(12.02), Inches(1.13), NAVY, NAVY, radius=True)
    add_text(
        slide,
        "poetry run python scripts/00_download_data.py",
        Inches(1.03),
        Inches(4.90),
        Inches(7.0),
        Inches(0.42),
        size=19,
        color=WHITE,
        bold=True,
        font=FONT_MONO,
    )
    add_text(
        slide,
        "safe to resume · exact-size validation",
        Inches(8.65),
        Inches(4.96),
        Inches(3.45),
        Inches(0.30),
        size=12,
        color=RGBColor(191, 207, 218),
        align=PP_ALIGN.RIGHT,
    )
    add_takeaway(
        slide,
        "Goal: follow every transformation using a complete recording and its synchronized physiology.",
        Inches(0.64),
        Inches(6.02),
        Inches(12.02),
    )


def build_raw_traces(prs):
    slide = add_slide(
        prs,
        "01 · Inspect the recording before building a graph",
        eyebrow="Track A · data literacy",
        step="script 01",
    )
    add_picture_contain(
        slide,
        FIGURES / "01_raw_dff_traces.png",
        Inches(0.63),
        Inches(1.27),
        Inches(8.42),
        Inches(5.37),
    )
    add_script_badge(slide, "01", "inspect_data.py", Inches(9.38), Inches(1.38))
    add_rich_lines(
        slide,
        [
            ("Why start here?", True, NAVY),
            ("See the measured fluorescence before using processed activity.", False, INK),
            ("100 random neurons", True, TEAL),
            ("full 40.7-minute session", False, INK),
            ("raw ΔF/F; offset only for display", False, INK),
        ],
        Inches(9.38),
        Inches(2.15),
        Inches(3.16),
        Inches(2.6),
        size=16,
        gap=9,
    )
    add_takeaway(
        slide,
        "Display choices do not change the traces used by later analyses.",
        Inches(9.28),
        Inches(5.25),
        Inches(3.35),
        color=TEAL,
    )
    add_caption(
        slide,
        "Output: results/figures/01_raw_dff_traces.png",
        Inches(0.72),
        Inches(6.68),
        Inches(7.0),
    )


def build_spatial_coverage(prs):
    slide = add_slide(
        prs,
        "01 · Check which cortical areas were sampled",
        eyebrow="Track A · spatial coverage",
        step="script 01",
    )
    add_picture_contain(
        slide,
        FIGURES / "01_spatial_distribution_by_area.png",
        Inches(0.60),
        Inches(1.20),
        Inches(10.02),
        Inches(5.68),
    )
    add_rich_lines(
        slide,
        [
            ("Each dot", True, TEAL),
            ("one recorded neuron", False, INK),
            ("Each color", True, TEAL),
            ("one Allen-atlas cortical area", False, INK),
            ("Why check?", True, CORAL),
            ("sampling determines which spatial claims the data can support", False, INK),
        ],
        Inches(10.85),
        Inches(1.52),
        Inches(1.82),
        Inches(3.90),
        size=12,
        gap=7,
    )
    add_takeaway(
        slide,
        "The map includes all recorded neurons, not only the later analysis subset.",
        Inches(10.73),
        Inches(5.58),
        Inches(1.98),
        color=TEAL,
        height=1.05,
        size=9,
    )


def build_physiology(prs):
    slide = add_slide(
        prs,
        "01 · Examine the physiological basis of state labels",
        eyebrow="Track A · EEG / EMG / behavior",
        step="script 01",
    )
    add_picture_contain(
        slide,
        FIGURES / "01_eeg_emg_state_classification.png",
        Inches(0.56),
        Inches(1.23),
        Inches(9.37),
        Inches(5.63),
    )
    add_pill(slide, "ALL NEURONS", Inches(10.25), Inches(1.42), Inches(1.43), color=NAVY)
    add_rich_lines(
        slide,
        [
            ("Raster", True, NAVY),
            ("positive deconvolution samples", False, INK),
            ("EEG spectrogram", True, NAVY),
            ("0.5–25 Hz power", False, INK),
            ("EMG envelope", True, NAVY),
            ("20–200 Hz RMS", False, INK),
            ("Deposited labels", True, NAVY),
            ("Wake · quiet wake · NREM · REM", False, INK),
        ],
        Inches(10.17),
        Inches(2.05),
        Inches(2.45),
        Inches(3.30),
        size=13,
        gap=4,
    )
    add_takeaway(
        slide,
        "Inspect the labels; this tutorial does not refit the state classifier.",
        Inches(10.12),
        Inches(5.55),
        Inches(2.55),
        color=CORAL,
        height=0.95,
        size=11,
    )


def build_connectivity(prs):
    slide = add_slide(
        prs,
        "02 · Build matched-density functional networks",
        eyebrow="Track A · functional connectivity",
        step="script 02",
    )
    add_picture_contain(
        slide,
        FIGURES / "02_connectivity_matrices.png",
        Inches(0.60),
        Inches(1.27),
        Inches(7.72),
        Inches(4.82),
    )
    add_picture_contain(
        slide,
        FIGURES / "02_fixed_density_graphs.png",
        Inches(8.57),
        Inches(1.27),
        Inches(4.12),
        Inches(3.16),
    )
    add_caption(
        slide,
        "Pearson correlation matrices",
        Inches(0.75),
        Inches(6.12),
        Inches(3.5),
    )
    add_caption(
        slide,
        "Strongest |r| edges at K = 5%",
        Inches(8.75),
        Inches(4.47),
        Inches(3.5),
    )
    add_takeaway(
        slide,
        "Match K across states so every graph has the same number of edges.",
        Inches(8.57),
        Inches(5.00),
        Inches(4.12),
        color=BLUE,
    )
    add_text(
        slide,
        "correlation ≠ anatomical connection",
        Inches(8.72),
        Inches(6.03),
        Inches(3.75),
        Inches(0.34),
        size=13,
        color=CORAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_modularity_pipeline(prs):
    slide = add_slide(
        prs,
        "03 · Detect modules with repeated Louvain optimization",
        eyebrow="Track A · modularity analysis",
        step="script 03",
    )
    add_picture_contain(
        slide,
        FIGURES / "03_modularity_pipeline_sleep.png",
        Inches(0.58),
        Inches(1.24),
        Inches(9.30),
        Inches(5.62),
    )
    add_rich_lines(
        slide,
        [
            ("Q asks:", True, TEAL),
            ("Are within-module edges more common than expected under the modularity null model?", False, INK),
            ("Louvain is stochastic", True, NAVY),
            ("Repeat the optimization; retain the maximum-Q partition.", False, INK),
            ("Labels are arbitrary", True, NAVY),
            ("Module 1 in one state is not automatically Module 1 in another.", False, INK),
        ],
        Inches(10.12),
        Inches(1.44),
        Inches(2.55),
        Inches(4.25),
        size=14,
        gap=7,
    )
    add_takeaway(
        slide,
        "Report the optimization settings: K, γ, number of runs, and neuron selection.",
        Inches(10.05),
        Inches(5.53),
        Inches(2.65),
        color=TEAL,
        height=1.05,
        size=11,
    )


def build_module_map(prs):
    slide = add_slide(
        prs,
        "03 · Functional modules can be spatially intermixed",
        eyebrow="Track A · interpret the partition",
        step="script 03",
    )
    add_picture_contain(
        slide,
        FIGURES / "03_spatial_modules_sleep.png",
        Inches(0.64),
        Inches(1.26),
        Inches(8.25),
        Inches(5.58),
    )
    add_text(
        slide,
        "Same field of view\nDifferent functional labels",
        Inches(9.36),
        Inches(1.67),
        Inches(3.05),
        Inches(0.96),
        size=24,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_rich_lines(
        slide,
        [
            ("Each dot", True, TEAL),
            ("one active neuron", False, INK),
            ("Each color", True, TEAL),
            ("one Louvain module", False, INK),
            ("Key observation", True, CORAL),
            ("functional membership is not simply a contiguous cortical parcel", False, INK),
        ],
        Inches(9.38),
        Inches(3.10),
        Inches(3.00),
        Inches(2.55),
        size=15,
        gap=8,
    )
    add_takeaway(
        slide,
        "Spatially intermixed does not mean spatially random—quantify it in script 06.",
        Inches(9.26),
        Inches(5.69),
        Inches(3.28),
        color=GOLD,
        height=0.90,
        size=11,
    )


def build_sample_compare(prs):
    slide = add_slide(
        prs,
        "04 · Compare Awake and NREM within mouse02",
        eyebrow="Track A · one-mouse result",
        step="script 04",
    )
    add_picture_contain(
        slide,
        FIGURES / "04_sample_state_comparison.png",
        Inches(0.62),
        Inches(1.27),
        Inches(8.62),
        Inches(5.38),
    )
    add_rich_lines(
        slide,
        [
            ("Repeat across", True, NAVY),
            ("complete stable-state windows", False, INK),
            ("graph densities K = 2%, 5%, 10%", False, INK),
            ("multiple Louvain seeds", False, INK),
            ("Interpret carefully", True, CORAL),
            ("windows from one animal are not independent mice", False, INK),
        ],
        Inches(9.62),
        Inches(1.52),
        Inches(2.75),
        Inches(3.30),
        size=15,
        gap=9,
    )
    add_takeaway(
        slide,
        "This is a complete worked example—not a population-level claim.",
        Inches(9.50),
        Inches(5.20),
        Inches(3.12),
        color=CORAL,
        height=0.90,
        size=11,
    )


def build_sample_scale(prs):
    slide = add_slide(
        prs,
        "05 · Rebuild networks at coarser spatial scales",
        eyebrow="Track A · scale dependence",
        step="script 05",
    )
    add_picture_contain(
        slide,
        FIGURES / "05_sample_coarse_grain_modularity.png",
        Inches(0.60),
        Inches(1.27),
        Inches(8.82),
        Inches(5.37),
    )
    add_rich_lines(
        slide,
        [
            ("At every scale", True, NAVY),
            ("group spatial neighbors", False, INK),
            ("average parcel activity", False, INK),
            ("recompute correlation", False, INK),
            ("threshold at fixed density", False, INK),
            ("rerun Louvain", False, INK),
        ],
        Inches(9.74),
        Inches(1.55),
        Inches(2.6),
        Inches(2.8),
        size=15,
        gap=9,
    )
    add_takeaway(
        slide,
        "One mouse can differ from the cohort pattern. Track B tests robustness across animals.",
        Inches(9.50),
        Inches(4.85),
        Inches(3.12),
        color=GOLD,
        height=1.15,
        size=11,
    )


def build_sample_spatial(prs):
    slide = add_slide(
        prs,
        "06 · Quantify how module geography changes with scale",
        eyebrow="Track A · spatial distribution",
        step="script 06",
    )
    add_picture_contain(
        slide,
        FIGURES / "06_sample_module_spatial_distribution.png",
        Inches(0.57),
        Inches(1.22),
        Inches(10.15),
        Inches(5.66),
    )
    add_rich_lines(
        slide,
        [
            ("Map", True, TEAL),
            ("module labels in cortical space", False, INK),
            ("Measure", True, TEAL),
            ("P(same module | distance)", False, INK),
            ("Compare", True, TEAL),
            ("single neurons vs 40-neuron parcels", False, INK),
        ],
        Inches(10.92),
        Inches(1.56),
        Inches(1.82),
        Inches(3.25),
        size=13,
        gap=8,
    )
    add_takeaway(
        slide,
        "Intermixed single-cell modules become more spatially localized after coarse-graining.",
        Inches(10.78),
        Inches(4.95),
        Inches(1.98),
        color=TEAL,
        height=1.38,
        size=9,
    )


def build_movie(prs):
    slide = add_slide(
        prs,
        "Supplemental · Follow the module map across spatial scales",
        eyebrow="Track A · multiscale synthesis",
        step="supplemental",
    )
    add_picture_contain(
        slide,
        figure_file(
            "multiscale_overview_mouse02_sleep_nrem.png",
            "07_multiscale_overview_mouse02_sleep_nrem.png",
        ),
        Inches(0.58),
        Inches(1.26),
        Inches(9.67),
        Inches(5.48),
    )
    add_pill(slide, "MP4 OUTPUT", Inches(10.55), Inches(1.52), Inches(1.50), color=CORAL)
    add_text(
        slide,
        "1 → 40\nneurons / parcel",
        Inches(10.53),
        Inches(2.17),
        Inches(1.66),
        Inches(0.92),
        size=23,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_DISPLAY,
    )
    add_rich_lines(
        slide,
        [
            ("Seven scales", True, TEAL),
            ("new network at each scale", False, INK),
            ("Color alignment", True, TEAL),
            ("keeps neighboring frames visually coherent", False, INK),
        ],
        Inches(10.50),
        Inches(3.49),
        Inches(1.82),
        Inches(1.75),
        size=13,
        gap=7,
    )
    add_text(
        slide,
        "results/movies/\nmultiscale_modules_…mp4",
        Inches(10.40),
        Inches(5.76),
        Inches(2.05),
        Inches(0.58),
        size=10,
        color=SLATE,
        font=FONT_MONO,
        align=PP_ALIGN.CENTER,
    )


def build_track_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.17), SLIDE_H, GOLD)
    add_text(
        slide,
        "TRACK B · OPTIONAL RESEARCH EXTENSION",
        Inches(0.84),
        Inches(0.75),
        Inches(6.8),
        Inches(0.30),
        size=11,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Now ask whether the result\nis robust across mice",
        Inches(0.82),
        Inches(1.38),
        Inches(7.7),
        Inches(1.45),
        size=34,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_text(
        slide,
        "poetry run python scripts/00_download_data.py --all",
        Inches(0.84),
        Inches(3.48),
        Inches(7.55),
        Inches(0.48),
        size=18,
        color=WHITE,
        bold=True,
        font=FONT_MONO,
    )
    add_text(
        slide,
        "10 calcium recordings · corresponding physiology · biological-mouse summaries",
        Inches(0.84),
        Inches(4.17),
        Inches(7.25),
        Inches(0.40),
        size=15,
        color=RGBColor(194, 210, 221),
    )
    add_rect(slide, Inches(0.84), Inches(5.02), Inches(7.06), Inches(1.13), RGBColor(27, 52, 74), RGBColor(42, 71, 94), radius=True)
    add_text(
        slide,
        "Defaults = exploratory course run\nPAPER_MODE = True = full neurons + 200 Louvain runs",
        Inches(1.14),
        Inches(5.30),
        Inches(6.45),
        Inches(0.65),
        size=15,
        color=WHITE,
        bold=True,
    )
    for index, (number, label, color) in enumerate(
        [
            ("07", "modularity across mice", TEAL),
            ("08", "scale dependence across mice", GOLD),
            ("09", "module geography across mice", CORAL),
        ]
    ):
        y = Inches(1.24 + index * 1.55)
        add_rect(slide, Inches(9.18), y, Inches(0.72), Inches(0.72), color, color, radius=True)
        add_text(
            slide,
            number,
            Inches(9.18),
            y + Inches(0.08),
            Inches(0.72),
            Inches(0.44),
            size=15,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            Inches(10.20),
            y + Inches(0.15),
            Inches(2.34),
            Inches(0.45),
            size=16,
            color=WHITE,
            bold=True,
        )
    add_text(
        slide,
        "Mouse 4 has two sleep days; average days within mouse before cohort inference.",
        Inches(9.12),
        Inches(6.12),
        Inches(3.35),
        Inches(0.62),
        size=12,
        color=RGBColor(194, 210, 221),
        align=PP_ALIGN.CENTER,
    )


def build_all_mice(prs):
    slide = add_slide(
        prs,
        "07 · Reproduce the state comparison across mice",
        eyebrow="Track B · population robustness",
        step="script 07",
    )
    add_picture_contain(
        slide,
        FIGURES / "07_all_mice_state_comparison.png",
        Inches(0.58),
        Inches(1.25),
        Inches(7.92),
        Inches(5.50),
    )
    add_picture_contain(
        slide,
        FIGURES / "07_modularity_per_mouse_sleep.png",
        Inches(8.73),
        Inches(1.25),
        Inches(3.96),
        Inches(2.40),
    )
    add_picture_contain(
        slide,
        FIGURES / "07_modularity_per_mouse_ane.png",
        Inches(8.73),
        Inches(3.91),
        Inches(3.96),
        Inches(2.40),
    )
    add_takeaway(
        slide,
        "The biological replicate is the mouse—not the window and not the recording day.",
        Inches(8.69),
        Inches(6.29),
        Inches(4.03),
        color=CORAL,
        height=0.66,
        size=11,
    )


def build_all_scale(prs):
    slide = add_slide(
        prs,
        "08 · Compare brain states across spatial scales",
        eyebrow="Track B · population scale analysis",
        step="script 08",
    )
    add_picture_contain(
        slide,
        FIGURES / "08_all_mice_coarse_grain_modularity.png",
        Inches(0.60),
        Inches(1.27),
        Inches(7.88),
        Inches(5.34),
    )
    add_picture_contain(
        slide,
        FIGURES / "08_all_mice_coarse_grain_modules.png",
        Inches(8.72),
        Inches(1.27),
        Inches(3.97),
        Inches(3.05),
    )
    add_rich_lines(
        slide,
        [
            ("Why all mice?", True, GOLD),
            ("One recording need not match the cohort.", False, INK),
            ("Read the uncertainty", True, NAVY),
            ("A CI crossing zero does not establish equivalence.", False, INK),
        ],
        Inches(8.83),
        Inches(4.50),
        Inches(3.67),
        Inches(1.55),
        size=14,
        gap=6,
    )
    add_takeaway(
        slide,
        "Scale changes both the number of nodes and the apparent module organization.",
        Inches(8.70),
        Inches(6.23),
        Inches(3.98),
        color=GOLD,
    )


def build_all_spatial(prs):
    slide = add_slide(
        prs,
        "09 · Reproduce module maps and distance profiles",
        eyebrow="Track B · population spatial analysis",
        step="script 09",
    )
    add_picture_contain(
        slide,
        FIGURES / "09_all_mice_module_maps.png",
        Inches(0.57),
        Inches(1.24),
        Inches(7.50),
        Inches(5.48),
    )
    add_picture_contain(
        slide,
        FIGURES / "09_all_mice_same_module_vs_distance.png",
        Inches(8.30),
        Inches(1.24),
        Inches(4.39),
        Inches(4.25),
    )
    add_rich_lines(
        slide,
        [
            ("Maps", True, TEAL),
            ("qualitative module geography", False, INK),
            ("Distance curves", True, TEAL),
            ("quantitative spatial organization", False, INK),
        ],
        Inches(8.48),
        Inches(5.67),
        Inches(2.30),
        Inches(0.90),
        size=13,
        gap=5,
    )
    add_takeaway(
        slide,
        "Use both representations: a compelling map is not a substitute for a summary across mice.",
        Inches(10.95),
        Inches(5.34),
        Inches(1.76),
        color=TEAL,
        height=1.34,
        size=8,
    )


def build_summary(prs):
    slide = add_slide(prs, "From tutorial to research project", eyebrow="Take-home workflow")
    steps = [
        ("1", "Inspect", "activity + EEG/EMG + state", TEAL),
        ("2", "Construct", "correlation + fixed-density graph", BLUE),
        ("3", "Partition", "repeated Louvain + modularity Q", CORAL),
        ("4", "Validate", "mice + densities + spatial scales", GOLD),
    ]
    for index, (number, title, body, color) in enumerate(steps):
        x = Inches(0.62 + index * 3.08)
        add_rect(slide, x, Inches(1.55), Inches(2.77), Inches(2.12), PALE, LINE, radius=True)
        add_rect(slide, x + Inches(0.22), Inches(1.78), Inches(0.48), Inches(0.48), color, color, radius=True)
        add_text(
            slide,
            number,
            x + Inches(0.22),
            Inches(1.80),
            Inches(0.48),
            Inches(0.32),
            size=13,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + Inches(0.85),
            Inches(1.80),
            Inches(1.55),
            Inches(0.34),
            size=18,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + Inches(0.22),
            Inches(2.55),
            Inches(2.30),
            Inches(0.66),
            size=14,
            color=SLATE,
            align=PP_ALIGN.CENTER,
        )

    add_text(
        slide,
        "Possible project directions",
        Inches(0.64),
        Inches(4.18),
        Inches(4.0),
        Inches(0.40),
        size=22,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    projects = [
        "module stability across time",
        "density and resolution sensitivity",
        "region-specific or layer-specific modules",
        "relationships to EEG spectral features",
        "alternative community-detection methods",
        "cross-mouse reproducibility",
    ]
    for index, project in enumerate(projects):
        column = index % 3
        row = index // 3
        x = Inches(0.66 + column * 4.07)
        y = Inches(4.83 + row * 0.75)
        add_rect(slide, x, y, Inches(3.72), Inches(0.52), WHITE, LINE, radius=True)
        add_rect(slide, x + Inches(0.15), y + Inches(0.16), Inches(0.16), Inches(0.16), TEAL, TEAL, radius=True)
        add_text(
            slide,
            project,
            x + Inches(0.43),
            y + Inches(0.11),
            Inches(3.08),
            Inches(0.27),
            size=12,
            color=INK,
            bold=True,
        )
    add_takeaway(
        slide,
        "Begin with scripts 00–06. Use 07–09 before making claims about robustness across mice.",
        Inches(0.65),
        Inches(6.25),
        Inches(12.03),
        color=TEAL,
        height=0.68,
    )


def validate_presentation(prs):
    if len(prs.slides) != 18:
        raise AssertionError(f"Expected 18 slides, found {len(prs.slides)}")
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise AssertionError(f"Slide {slide_number}: shape starts outside slide")
            if shape.left + shape.width > SLIDE_W + Inches(0.01):
                raise AssertionError(f"Slide {slide_number}: shape extends past right edge")
            if shape.top + shape.height > SLIDE_H + Inches(0.01):
                raise AssertionError(f"Slide {slide_number}: shape extends past bottom edge")


def main():
    require_figures()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "CSHA functional-network modularity hands-on"
    prs.core_properties.subject = "One-recording tutorial and all-mice research extension"
    prs.core_properties.author = "CSHA Neural Data Science course"
    prs.core_properties.keywords = "calcium imaging, EEG, EMG, modularity, Louvain, network analysis"

    build_title(prs)
    build_roadmap(prs)
    build_download(prs)
    build_raw_traces(prs)
    build_spatial_coverage(prs)
    build_physiology(prs)
    build_connectivity(prs)
    build_modularity_pipeline(prs)
    build_module_map(prs)
    build_sample_compare(prs)
    build_sample_scale(prs)
    build_sample_spatial(prs)
    build_movie(prs)
    build_track_b(prs)
    build_all_mice(prs)
    build_all_scale(prs)
    build_all_spatial(prs)
    build_summary(prs)

    validate_presentation(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    reopened = Presentation(OUTPUT)
    validate_presentation(reopened)
    print(f"saved -> {OUTPUT}")
    print(f"slides: {len(reopened.slides)}")
    print(f"size: {OUTPUT.stat().st_size / 1e6:.2f} MB")


if __name__ == "__main__":
    main()
